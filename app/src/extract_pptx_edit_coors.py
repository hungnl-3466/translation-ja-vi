from pptx import Presentation
import matplotlib.pyplot as plt
import matplotlib.patches as patches

# Chuyển đổi từ EMU sang inches
EMU_TO_INCH = 914400

def extract_text_and_coordinates(pptx_path):
    prs = Presentation(pptx_path)
    slides_data = []

    for slide in prs.slides:
        slide_data = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    left = shape.left / EMU_TO_INCH
                    top = shape.top / EMU_TO_INCH
                    width = shape.width / EMU_TO_INCH
                    height = shape.height / EMU_TO_INCH
                    slide_data.append((text, left, top, width, height))
        slides_data.append(slide_data)
    return slides_data

def draw_text_from_pptx(slides_data):
    for i, slide_data in enumerate(slides_data):
        fig, ax = plt.subplots()
        ax.set_xlim(0, 13.33)  # Mặc định slide rộng 13.33 inches
        ax.set_ylim(7.5, 0)    # Mặc định slide cao 7.5 inches, đảo ngược trục y
        
        for text, left, top, width, height in slide_data:
            # Vẽ text
            plt.text(left, top, text, fontsize=12, verticalalignment='top', wrap=True, bbox=dict(facecolor='white', alpha=0.5))

            # Vẽ khung chữ nhật bao quanh text
            rect = patches.Rectangle((left, top), width, height, linewidth=1, edgecolor='r', facecolor='none')
            ax.add_patch(rect)

        plt.title(f"Slide {i+1}")
        plt.show()

# Đường dẫn tới file PPTX của bạn
pptx_path = "Nhom8_QLDA_Slide.pptx"
slides_data = extract_text_and_coordinates(pptx_path)
draw_text_from_pptx(slides_data)
