from pptx import Presentation
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle

def get_slide_size(presentation):
    # Lấy kích thước slide từ cài đặt của presentation
    slides_width = presentation.slide_width
    slides_height = presentation.slide_height
    return slides_width, slides_height

def draw_text_boxes_with_bounding_boxes(pptx_path):
    prs = Presentation(pptx_path)

    for slide_idx, slide in enumerate(prs.slides):
        fig, ax = plt.subplots()
        ax.imshow([[1, 1], [1, 1]])  # Tạo một hình ảnh trắng để vẽ bounding box

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            rect = Rectangle((left, top), width, height, linewidth=1, edgecolor='r', facecolor='none')
            ax.add_patch(rect)
            ax.text(left, top, text, fontsize=8, verticalalignment='top', color='r')

        ax.axis('off')
        slides_width, slides_height = get_slide_size(prs)
        ax.set_xlim(0, slides_width)
        ax.set_ylim(slides_height, 0)  # Lật ngược trục y để phù hợp với hệ tọa độ của PowerPoint
        plt.title(f"Slide {slide_idx + 1}")
        plt.show()


pptx_path = 'Nhom8_QLDA_Slide.pptx'
draw_text_boxes_with_bounding_boxes(pptx_path)