from pptx import Presentation
import cv2
import numpy as np
from PIL import ImageFont, ImageDraw, Image
import json

EMU_TO_INCH = 914400
EMU_TO_PIXEL = EMU_TO_INCH / 96  # Chuyển đổi từ EMU sang pixel với DPI = 96

def extract_text_and_coordinates(pptx_path):
    prs = Presentation(pptx_path)
    slides_data = []

    for slide in prs.slides:
        slide_data = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text_content = []
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    if text.strip():
                        text_content.append(text)
            
            if text_content:
                text = "\n".join(text_content)
                left = int(shape.left / EMU_TO_PIXEL)
                top = int(shape.top / EMU_TO_PIXEL)
                width = int(shape.width / EMU_TO_PIXEL)
                height = int(shape.height / EMU_TO_PIXEL)
                slide_data.append((text, left, top, width, height))
        slides_data.append(slide_data)
    return slides_data

def calculate_text_width(text, font):
    image = Image.new('RGB', (1000, 100), (255, 255, 255))
    draw = ImageDraw.Draw(image)
    bbox = draw.textbbox((10, 10), text, font=font)
    text_width, text_height = bbox[2] - bbox[0], bbox[3] - bbox[1]
    return text_width

def cut_text_by_pixel_length(text, font_path, font_size, max_pixel_length):
    font = ImageFont.truetype(font_path, font_size)
    current_length = 0
    cut_text = ""
    
    for char in text:
        char_length = calculate_text_width(char, font)
        if current_length + char_length > max_pixel_length:
            break
        cut_text += char
        current_length += char_length
    
    return cut_text


def draw_text_from_pptx(slides_data, font_path):
    
    dict_slide = {}
    
    dict_text = {
        "id": int,
        "bbox": [],
        "text": str
    }
    
    font_size = 24
    dpi = 96
    slide_width = int(13.33 * dpi)
    slide_height = int(7.5 * dpi)

    for idx_slide, slide_data in enumerate(slides_data):
        name_slide = "slide_{}".format(idx_slide)
        print(name_slide)
        dict_slide[name_slide] = []
        # Tạo ảnh trắng kích thước slide
        slide_img = np.ones((slide_height, slide_width, 3), dtype=np.uint8) * 255

        # Chuyển đổi ảnh sang PIL Image
        pil_img = Image.fromarray(slide_img)
        print("========= ==========================")
        idx_dict_text = 0
        for text, left, top, width, height in slide_data:
            draw = ImageDraw.Draw(pil_img)
            font = ImageFont.truetype(font_path, font_size)
            # Vẽ văn bản sử dụng PIL
            bbox = draw.textbbox((left, top), text, font=font)
            text_width, text_height = bbox[2] - bbox[0], bbox[3] - bbox[1]

            text_temp = text
            # TO DO: rule length of text to down the line
            length_text_line = left + text_width
            length_bbox = left + width
  
            temp_length_text_line = length_text_line
            temp_top = top
            lines_list = []

            dict_text = {
                "id": idx_dict_text,
                "bbox": [top, left, length_bbox, top + height],
                "text": text
            }
            
            dict_slide[name_slide].append(dict_text)
            
            # while length_text_line > temp_length_bbox:
            if temp_length_text_line > length_bbox:
                num_lines = temp_length_text_line // width
                for i in range(num_lines):
                    cut_text = cut_text_by_pixel_length(text_temp, font_path, font_size, width)
                    print(cut_text)
                    idx = len(cut_text)
                    text_temp = text_temp[idx:]
                    temp_length_text_line = temp_length_text_line - width
                    lines_list.append(cut_text)
                    
                    draw.text((left, temp_top), cut_text, font=font, fill=(0, 0, 0))
                    temp_top = temp_top + text_height
            
            else:
                draw.text((left, top), text, font=font, fill=(0, 0, 0))
                # Đo kích thước văn bản

            
            
            slide_img = np.array(pil_img)  # Chuyển đổi lại sang ảnh OpenCV
            
            # # // TO DO: Draw start point of the texttop left
            # cv2.circle(slide_img, (left, top), 1, (0,255,0), 10)
            # # // TO DO: Draw length text point
            # cv2.circle(slide_img, (left + text_width, top), 1, (0,255,255), 10)
            # # // TO DO: Draw end point of the text 
            # cv2.circle(slide_img, (length_bbox, top), 1, (255,255,0), 10)
            # // TO DO: Draw bound boxes are extracted with structure from pptx            
            cv2.rectangle(slide_img, (left, top), (length_bbox, top + height), (0, 0, 255), 1)
            pil_img = Image.fromarray(slide_img)  # Chuyển đổi lại sang PIL để vẽ văn bản tiếp theo
            idx_dict_text += 1
            print("============================================== ")
        # Hiển thị slide sử dụng OpenCV
        slide_img = np.array(pil_img)
        cv2.imshow(f"Slide {idx_slide+1}", slide_img)
        # idx_slide += 1
        cv2.waitKey(1)
        cv2.destroyAllWindows()
    
    print(dict_slide)
    with open('dict_slide_text.json', 'w', encoding='utf-8') as f:
        json.dump(dict_slide, f, ensure_ascii=False, indent=4)
    
# Đường dẫn tới file PPTX và phông chữ tiếng Nhật của bạn
pptx_path = "Ja_ver_Sun_AI_Development.pptx"
font_path = 'NotoSansJP-VariableFont_wght.ttf'

slides_data = extract_text_and_coordinates(pptx_path)
draw_text_from_pptx(slides_data, font_path)
