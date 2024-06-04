from pptx import Presentation

def extract_text_and_coordinates(pptx_path):
    prs = Presentation(pptx_path)
    slides_data = []

    for slide in prs.slides:
        slide_data = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    slide_data.append((text, left, top, width, height))
        slides_data.append(slide_data)
    return slides_data

pptx_path = "Nhom8_QLDA_Slide.pptx"
slides_data = extract_text_and_coordinates(pptx_path)

# Kiểm tra dữ liệu đã trích xuất
for i, slide_data in enumerate(slides_data):
    print(f"Slide {i+1}:")
    for text, left, top, width, height in slide_data:
        print(f"Text: {text}, Left: {left}, Top: {top}, Width: {width}, Height: {height}")


import matplotlib.pyplot as plt
import matplotlib.patches as patches

def draw_text_from_pptx(slides_data):
    for i, slide_data in enumerate(slides_data):
        fig, ax = plt.subplots()
        ax.set_xlim(0, 10)  # Tùy chỉnh giới hạn x
        ax.set_ylim(0, 10)  # Tùy chỉnh giới hạn y
        ax.invert_yaxis()  # Đảo ngược trục y để phù hợp với tọa độ PPTX

        for text, left, top, width, height in slide_data:
            # Chuyển đổi tọa độ từ EMU sang inches (nếu cần)
            left_inch = left / 914400
            top_inch = top / 914400
            width_inch = width / 914400
            height_inch = height / 914400

            # Vẽ text
            plt.text(left_inch, top_inch, text, fontsize=12, verticalalignment='top', bbox=dict(facecolor='white', alpha=0.5))

            # Vẽ khung chữ nhật bao quanh text
            rect = patches.Rectangle((left_inch, top_inch), width_inch, height_inch, linewidth=1, edgecolor='r', facecolor='none')
            ax.add_patch(rect)

        plt.title(f"Slide {i+1}")
        plt.show()

draw_text_from_pptx(slides_data)





            
            
# 'Nhom8_QLDA_Slide.pptx'