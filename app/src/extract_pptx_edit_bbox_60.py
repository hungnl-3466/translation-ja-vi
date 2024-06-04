from pptx import Presentation
import matplotlib.pyplot as plt
import matplotlib.patches as patches

EMU_TO_INCH = 914400

def extract_text_and_coordinates(pptx_path):
    prs = Presentation(pptx_path)
    slides_data = []

    for slide in prs.slides:
        slide_data = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text_content = []
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    if text.strip():
                        text_content.append(text)
            
            if text_content:
                text = "\n".join(text_content)
                left = shape.left / EMU_TO_INCH
                top = shape.top / EMU_TO_INCH
                width = shape.width / EMU_TO_INCH
                height = shape.height / EMU_TO_INCH
                slide_data.append((text, left, top, width, height))
        slides_data.append(slide_data)
    return slides_data

pptx_path = "Ja_ver_Sun_AI_Development.pptx"
slides_data = extract_text_and_coordinates(pptx_path)

# In ra văn bản để kiểm tra
for slide_data in slides_data:
    for text, left, top, width, height in slide_data:
        print(f"Text: {text}, Coordinates: ({left}, {top}, {width}, {height})")

import matplotlib.font_manager as fm

# Đường dẫn tới file phông chữ Noto Sans CJK JP
font_path = 'NotoSansJP-VariableFont_wght.ttf'
font_prop = fm.FontProperties(fname=font_path)

def draw_text_from_pptx(slides_data):
    print("=========================================================")
    for i, slide_data in enumerate(slides_data):
        fig, ax = plt.subplots(figsize=(13.33, 7.5))
        ax.set_xlim(0, 13.33)  # Slide width in inches
        ax.set_ylim(7.5, 0)    # Slide height in inches, inverted y-axis
        print("-------------------------")
        for text, left, top, width, height in slide_data:
            # Vẽ text với phông chữ hỗ trợ tiếng Nhật
            print("--- Coordinates----")
            print(left, top)
            print("-----Text-----")
            print(text)
            ax.text(left, top, text, fontsize=12, fontproperties=font_prop, verticalalignment='top', wrap=True, bbox=dict(facecolor='white', alpha=0.5))
            
            # Vẽ khung chữ nhật bao quanh text
            rect = patches.Rectangle((left, top), width, height, linewidth=1, edgecolor='r', facecolor='none')
            ax.add_patch(rect)

        plt.title(f"Slide {i+1}")
        plt.axis('off')
        plt.show()

draw_text_from_pptx(slides_data)
# Đường dẫn tới file PPTX của bạn
# pptx_path = "Ja_ver_Sun_AI_Development.pptx"
# slides_data = extract_text_and_coordinates(pptx_path)
draw_text_from_pptx(slides_data)
