from pptx import Presentation
import json


def replace_all_texts(pptx_path, output_path, json_path):
    # Mở file PowerPoint
    prs = Presentation(pptx_path)
    
    # Đọc dữ liệu từ file JSON
    with open(json_path, 'r', encoding='utf-8') as f:
        data_batch = json.load(f)
    
    # Duyệt qua tất cả các slide
    slide_cnt = 0
    for slide in prs.slides:
        name_slide = f"slide_{slide_cnt}"
        
        if name_slide in data_batch:
            dict_slide = data_batch[name_slide]
            id_cnt = 0

            # Duyệt qua tất cả các shape trong mỗi slide
            for shape in slide.shapes:
                if id_cnt >= len(dict_slide):
                    break  # Không còn văn bản để thay thế

                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    if shape.text_frame.text == '':
                        continue
                    else:
                        dict_text = dict_slide[id_cnt]
                        new_text = dict_text['text']

                        # Thay thế văn bản mà không làm thay đổi định dạng gốc
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.runs:
                                # Sử dụng một `run` đầu tiên để đặt toàn bộ văn bản mới vào đó, các `run` còn lại sẽ bị xoá đi
                                first_run = paragraph.runs[0]
                                first_run.text = new_text
                                for run in paragraph.runs[1:]:
                                    run.text = ''

                        id_cnt += 1
                
                # Kiểm tra nếu shape là bảng
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if id_cnt >= len(dict_slide):
                                break  # Không còn văn bản để thay thế

                            if cell.text == '':
                                continue
                            else:
                                dict_text = dict_slide[id_cnt]
                                new_text = dict_text['text']
                                
                                # Thay thế văn bản mà không làm thay đổi định dạng gốc
                                for paragraph in cell.text_frame.paragraphs:
                                    if paragraph.runs:
                                        # Sử dụng một `run` đầu tiên để đặt toàn bộ văn bản mới vào đó, các `run` còn lại sẽ bị xoá đi
                                        first_run = paragraph.runs[0]
                                        first_run.text = new_text
                                        for run in paragraph.runs[1:]:
                                            run.text = ''

                                id_cnt += 1
        slide_cnt += 1

    # Lưu file PowerPoint mới
    prs.save(output_path)



pptx_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/app/src/Ja_ver_Sun_AI_Development.pptx'
output_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/app/src/test_3.pptx'
json_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/json_data/dict_slide_text_trans.json'

# Gọi hàm với đường dẫn file PowerPoint, file xuất và file JSON
replace_all_texts(pptx_path, output_path, json_path)

