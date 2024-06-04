from pptx import Presentation
import json
from pptx.util import Pt

def replace_all_texts(pptx_path, output_path, json_path):
    # Mở file PowerPoint
    prs = Presentation(pptx_path)
    
    # Đọc dữ liệu từ file json
    with open(json_path, 'r', encoding='utf-8') as f:
        data_batch = json.load(f)

    # Duyệt qua tất cả các slide
    slide_cnt = 0
    for slide in prs.slides:
        name_slide = "slide_{}".format(slide_cnt)
        print("[SLIDE NAME]", name_slide)
        dict_slide = data_batch[name_slide]
        id_cnt = 0

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                if shape.text_frame.text == '':
                    continue
                else:  
                    dict_text = dict_slide[id_cnt]
                    original_paragraphs = shape.text_frame.paragraphs
                    new_text = dict_text['text']
                    
                    # Duyệt qua từng đoạn văn bản và thay thế nội dung
                    for paragraph in original_paragraphs:
                        for run in paragraph.runs:
                            # Sao chép thuộc tính định dạng hiện tại
                            font = run.font
                            size = font.size
                            bold = font.bold
                            italic = font.italic
                            underline = font.underline
                            if font.color.type == 1:  # RGB color
                                color = font.color.rgb
                            else:
                                color = None
                            
                            # Thay thế nội dung của run
                            run.text = new_text
                            new_text = ''  # Đảm bảo rằng chỉ thay thế văn bản một lần

                            # Khôi phục thuộc tính định dạng
                            run.font.size = size
                            run.font.bold = bold
                            run.font.italic = italic
                            run.font.underline = underline
                            if color:
                                run.font.color.rgb = color

                    id_cnt += 1

            # Kiểm tra nếu shape là bảng
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text == '':
                            continue
                        else:
                            dict_text = dict_slide[id_cnt]
                            new_text = dict_text['text']
                            original_paragraphs = cell.text_frame.paragraphs

                            # Duyệt qua từng đoạn văn bản và thay thế nội dung
                            for paragraph in original_paragraphs:
                                for run in paragraph.runs:
                                    # Sao chép thuộc tính định dạng hiện tại
                                    font = run.font
                                    size = font.size
                                    bold = font.bold
                                    italic = font.italic
                                    underline = font.underline
                                    if font.color.type == 1:  # RGB color
                                        color = font.color.rgb
                                    else:
                                        color = None

                                    # Thay thế nội dung của run
                                    run.text = new_text
                                    new_text = ''  # Đảm bảo rằng chỉ thay thế văn bản một lần

                                    # Khôi phục thuộc tính định dạng
                                    run.font.size = size
                                    run.font.bold = bold
                                    run.font.italic = italic
                                    run.font.underline = underline
                                    if color:
                                        run.font.color.rgb = color

                            id_cnt += 1


        slide_cnt += 1

    # Lưu file PowerPoint mới
    prs.save(output_path)

# Đường dẫn tới file PowerPoint ban đầu và file đầu ra
# pptx_path = 'input.pptx'
# output_path = 'output.pptx'
# new_text = 'Text mới chèn vào'

# # Gọi hàm replace_all_texts
# replace_all_texts(pptx_path, output_path, new_text)


# Đường dẫn tới file PowerPoint ban đầu và file đầu ra
pptx_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/app/src/Ja_ver_Sun_AI_Development.pptx'
output_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/app/src/test_4.pptx'
json_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/json_data/dict_slide_text_trans.json'
# Gọi hàm clear_all_texts
replace_all_texts(pptx_path, output_path, json_path)
