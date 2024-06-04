from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import json

def replace_text_in_shape(shape, new_texts):
    """
    Replace text in a shape based on new_texts from JSON.
    """
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if new_texts:
                text_info = new_texts.pop(0)
                run.text = text_info["text"]
                if text_info.get("font_name"):
                    run.font.name = text_info["font_name"]
                if text_info.get("font_size"):
                    run.font.size = Pt(text_info["font_size"])
                if text_info.get("bold") is not None:
                    run.font.bold = text_info["bold"]
                if text_info.get("italic") is not None:
                    run.font.italic = text_info["italic"]
                if text_info.get("color"):
                    r, g, b = text_info["color"]
                    run.font.color.rgb = RGBColor(r, g, b)

def replace_text_in_table(table, new_texts):
    """
    Replace text in a table based on new_texts from JSON.
    """
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if new_texts:
                        text_info = new_texts.pop(0)
                        run.text = text_info["text"]
                        if text_info.get("font_name"):
                            run.font.name = text_info["font_name"]
                        if text_info.get("font_size"):
                            run.font.size = Pt(text_info["font_size"])
                        if text_info.get("bold") is not None:
                            run.font.bold = text_info["bold"]
                        if text_info.get("italic") is not None:
                            run.font.italic = text_info["italic"]
                        if text_info.get("color"):
                            r, g, b = text_info["color"]
                            run.font.color.rgb = RGBColor(r, g, b)

def parser_text(pptx_path, output_path, json_path):
    prs = Presentation(pptx_path)
    
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    for id_slide, slide in enumerate(prs.slides):
        name_slide = f"slide_{id_slide}"
        if name_slide in data:
            dict_slide = data[name_slide]
            for item in dict_slide:
                id_text = item['id']
                list_text = item['list_text']
                # if id_text < len(slide.shapes):
                shape = slide.shapes[id_text]
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    replace_text_in_shape(shape, list_text)
                elif shape.has_table:
                    replace_text_in_table(shape.table, list_text)
    
    prs.save(output_path)


pptx_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/app/src/Ja_ver_Sun_AI_Development.pptx'
json_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/app/aplications/data_test.json'
output_path = 'output.pptx'

# Usage
parser_text(pptx_path, output_path, json_path)