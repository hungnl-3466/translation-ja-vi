from pptx import Presentation

def replace_all_texts(pptx_path, output_path, new_text):
    # Mở file PowerPoint
    prs = Presentation(pptx_path)
    
    # Duyệt qua tất cả các slide
    slide_cnt = 0
    for slide in prs.slides:
        print("-------------------------------------")
        print("Slide: ", slide_cnt)
        # Duyệt qua tất cả các shape trong mỗi slide
        text_cnt = 0
        for shape in slide.shapes:
            print("Index text: ", text_cnt)
            # Kiểm tra nếu shape có chứa text
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                # Xóa text cũ và chèn text mới
                if shape.text_frame.text is None:
                    pass
                else:
                    print(shape.text_frame.text)
                    shape.text_frame.clear()  # Xóa toàn bộ text cũ
                    p = shape.text_frame.add_paragraph()
                    p.text = new_text
                    print("Text in shape")
            
            # Kiểm tra nếu shape là bảng
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text is None:
                            pass
                        else:
                            print(cell.text)
                            cell.text = new_text
                            print("text in shell")
                            
            # Kiểm tra nếu shape là placeholder và có chứa text
            # elif shape.is_placeholder:
            #     if hasattr(shape, "text_frame") and shape.text_frame is not None:
            #         shape.text_frame.clear()
            #         p = shape.text_frame.add_paragraph()
            #         print(p.text)
            #         p.text = new_text
            #         print("text in placeholder")
            
            
            text_cnt += 1
        slide_cnt += 1
    # Lưu file PowerPoint mới
    prs.save(output_path)

# Đường dẫn tới file PowerPoint ban đầu và file đầu ra
pptx_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/app/src/Ja_ver_Sun_AI_Development.pptx'
output_path = 'test.pptx'
new_text = 'Text mới chèn vào'

# Gọi hàm replace_all_texts
replace_all_texts(pptx_path, output_path, new_text)
