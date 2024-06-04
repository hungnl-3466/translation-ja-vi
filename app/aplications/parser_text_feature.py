from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import json
# from pptx.util import Pt



def parser_text(pptx_path, output_pat, json_path):
    
    prs = Presentation(pptx_path)
    
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    id_slide = 0
    
    for slide in prs.slides:
        name_slide = "slide_{}".format(id_slide)
        print("Name slide: ", name_slide)
        dict_slide = data[name_slide]        
        id_text = 0
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                if shape.text_frame.text == '':
                    continue
                else:
                    dict_list_text = dict_slide[id_text]
                    cnt = 0
                
                    idx_run = 0
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text.strip()
                            if text:
                                dict_text = \
                                    dict_list_text["list_text"][idx_run]
                                text_parser = dict_text["text"]
                                run.text = text_parser
                                idx_run+=1
                                cnt+=1
                    id_text += 1
                                       
            elif shape.has_table:
                table = shape.table
                # id_text = 0
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text == '':
                            continue
                        else:                
                            idx_run = 0
                            dict_list_text = dict_slide[id_text]
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text = run.text.strip()
                                    if text:
                                        dict_text = dict_list_text["list_text"][idx_run]
                                        text_parser = dict_text["text"]
                                        run.text = text_parser
                                        
                                        idx_run+=1
                            id_text += 1
            
        # if id_slide == 23:
        #     break
        id_slide += 1
    
    prs.save(output_path)

    
pptx_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/app/src/Ja_ver_Sun_AI_Development.pptx'
json_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/app/aplications/translated_dict.json'
output_path = 'output_ver_feature_json.pptx'

parser_text(pptx_path, output_path, json_path)
