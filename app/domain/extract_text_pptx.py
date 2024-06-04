from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import json



EMU_TO_PIXEL = 914400 / 96




dict_slide = {}
dict_text = {}


def extract_text(pptx_path):
    prs = Presentation(pptx_path)
    print("aaa")
    id_slide = 0
    for slide in prs.slides:

        name_slide = "slide_{}".format(id_slide)
        print("=================================")
        print(name_slide)
        dict_slide[name_slide] = []
        
        id_text = 0
        for shape in slide.shapes:
            # Extract text from text frames
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                if shape.text_frame.text == '':
                    continue
                else: 
                    print(id_text)
                    # dict_text = dict_slide[id_text]
                    
                    list_text = []
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            print("text texxxxtttttt: ", run.text)
                            text = run.text.strip()
                            if text:
                                font = run.font
                                font_name = font.name
                                font_size = font.size.pt if font.size else None
                                bold = font.bold
                                italic = font.italic
                                color = font.color.rgb if font.color and font.color.type == 1 else None
                                                            
                                list_text.append({
                                    'text': text,
                                    'font_name': font_name,
                                    'font_size': font_size,
                                    'bold': bold,
                                    'italic': italic,
                                    'color': color
                                })
                                # list_text.append(text)
                                # print(list_text)
                    print(list_text)  
                    
                    dict_text = {
                        "id": id_text,
                        "list_text": list_text
                    }
                    dict_slide[name_slide].append(dict_text)
                    id_text += 1
                    
            # Extract text from tables
            elif shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text == '':
                            continue
                        else:
                            list_text=[]
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text = run.text.strip()
                                    if text:
                                        font = run.font
                                        font_name = font.name
                                        font_size = font.size.pt if font.size else None
                                        bold = font.bold
                                        italic = font.italic
                                        color = font.color.rgb if font.color and font.color.type == 1 else None

                                        # list_text.append(text)
                                        list_text.append({
                                            'text': text,
                                            'font_name': font_name,
                                            'font_size': font_size,
                                            'bold': bold,
                                            'italic': italic,
                                            'color': color
                                        })
                                        # print(list_text)
                            print(list_text)  
                            dict_text = {
                                "id": id_text,
                                "list_text": list_text
                            }
                            id_text += 1

            
            
            # dict_slide[name_slide] = dict_text
                            dict_slide[name_slide].append(dict_text)
        id_slide += 1
                
        
    print(dict_slide)
    with open('extracted_text_feature.json', 'w', encoding='utf-8') as f:
        dict_slide_json = json.dump(dict_slide, f, ensure_ascii=False,indent = 4)
        # s
    # print(dict_slide_json)
        # slides_data.append(slide_data)
    # return slides_data
    
pptx_path = '/media/benu/DATA/sun-asterisk/translation-ja-vi/app/src/Ja_ver_Sun_AI_Development.pptx'

extract_text(pptx_path)
