
from pptx import Presentation
# import cv2
import numpy as np
# from PIL import ImageFont, ImageDraw, Image
import json
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

import sys
sys.path.insert(0, "app")

class Pineline():
    def __init__(self):
        self.EMU_TO_INCH = 914400
        self.EMU_TO_PIXEL = self.EMU_TO_INCH / 96
        self.font_path = "app/font/DejaVuSans.ttf"
        
    def extract_text(self, pptx_content):
        
        prs = Presentation(pptx_content)
        slides_data = []

        for slide in prs.slides:
            slide_data = []
            for shape in slide.shapes:
                # Extract text from text frames
                if shape.has_text_frame:
                    text_content = []
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text
                            if text.strip():
                                text_content.append(text)
                    if text_content:
                        text = "\n".join(text_content)
                        left = int(shape.left / self.EMU_TO_PIXEL)
                        top = int(shape.top / self.EMU_TO_PIXEL)
                        width = int(shape.width / self.EMU_TO_PIXEL)
                        height = int(shape.height / self.EMU_TO_PIXEL)
                        slide_data.append((text, left, top, width, height))

                # Extract text from tables
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    table_left = int(shape.left / self.EMU_TO_PIXEL)
                    table_top = int(shape.top / self.EMU_TO_PIXEL)
                    table_width = int(shape.width / self.EMU_TO_PIXEL)
                    table_height = int(shape.height / self.EMU_TO_PIXEL)

                    row_height = table_height / len(table.rows)
                    col_widths = [table_width / len(table.columns)] * len(table.columns)

                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            text = cell.text.strip()
                            if text:
                                left = table_left + int(col_idx * col_widths[col_idx])
                                top = table_top + int(row_idx * row_height)
                                width = int(col_widths[col_idx])
                                height = int(row_height)
                                slide_data.append((text, left, top, width, height))
            slides_data.append(slide_data)
        # print(len(slides_data))
        return slides_data

    def calculate_text_width(self, text, font):
        image = Image.new('RGB', (1000, 100), (255, 255, 255))
        draw = ImageDraw.Draw(image)
        bbox = draw.textbbox((10, 10), text, font=font)
        text_width, text_height = bbox[2] - bbox[0], bbox[3] - bbox[1]
        return text_width

    
    def cut_text_by_pixel_length(self, text, font_path, font_size, max_pixel_length):
        font = ImageFont.truetype(font_path, font_size)
        current_length = 0
        cut_text = ""
        
        for char in text:
            char_length = self.calculate_text_width(char, font)
            if current_length + char_length > max_pixel_length:
                break
            cut_text += char
            current_length += char_length
        
        return cut_text
    
    def draw_text_from_pptx(self, slides_data):
        
        dict_slide = {}
        
        dict_text = {
            "id": int,
            "bbox": [],
            "text": str
        }
        
        font_size = 24
        dpi = 96
        slide_width = int(13.33 * dpi)
        slide_height = int(7.5 * dpi)

        for idx_slide, slide_data in enumerate(slides_data):
            name_slide = "slide_{}".format(idx_slide)
            dict_slide[name_slide] = []
            # Tạo ảnh trắng kích thước slide
            slide_img = np.ones((slide_height, slide_width, 3), dtype=np.uint8) * 255

            # Chuyển đổi ảnh sang PIL Image
            pil_img = Image.fromarray(slide_img)
            # print("====================================")
            idx_dict_text = 0
            for text, left, top, width, height in slide_data:
                dict_text = {
                    "id": idx_dict_text,
                    "text": text
                }
                
                dict_slide[name_slide].append(dict_text)
                
            
                idx_dict_text += 1

        return dict_slide
    def split_batch(self, dict_output, num_of_batch:int):
        dict_batch = {}
        cnt = 0
        cnt_batch = 0
        batch = []
        for name_slide in dict_output:
            dict_slide = dict_output[name_slide]
            dict_temp = {}
            dict_temp[name_slide] = dict_slide
            batch.append(dict_temp)
            if cnt > 0 and cnt % num_of_batch == 0:
                name_batch = "batch_{}".format(cnt_batch)
                dict_batch[name_batch] = batch
                batch = []
                cnt_batch += 1
            if len(dict_output) - cnt < num_of_batch:
                name_last_batch = "batch_{}".format(cnt_batch+1)
                dict_batch[name_last_batch] = batch
            cnt += 1
            
        return dict_batch
    
    def parser_text(self, pptx_path):
        
        prs = Presentation(pptx_path)
    
        # Duyệt qua tất cả các slide
        slide_cnt = 0
        f = open('/media/benu/DATA/sun-asterisk/translation-ja-vi/json_data/dict_slide_text_trans.json')
        data_batch = json.load(f)
        # print(data_batch)
        # exit()
        for slide in prs.slides:
            name_slide = "slide_{}".format(slide_cnt)
            # Duyệt qua tất cả các shape trong mỗi slide
            print("==============================================================")
            print("[SLIDE NAME]", name_slide)
            # print(data_batch[name_slide])
            print("--------------")
                # Kiểm tra nếu shape có chứa text
            dict_slide = data_batch[name_slide]
            id_cnt = 0
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    # check = shape.text_frame.text
                    # print("frame text: ", check)
                    if shape.text_frame.text == '':
                        continue
                    else:
                    # if hasattr(shape, "text_frame") and shape.text_frame is not None:
                        dict_text = dict_slide[id_cnt]
                        # Xóa text cũ và chèn text mới
                        shape.text_frame.clear()  # Xóa toàn bộ text cũ
                        p = shape.text_frame.add_paragraph()
                        print(dict_text)
                        print("Parser in Place holder")
                        p.text =  dict_text['text'] 
                        p.font.size = Pt(10)
                        print(id_cnt)
                        id_cnt += 1
                    
                    
                # Kiểm tra nếu shape là bảng
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text == '':
                                continue
                            else:
                                print("Parser in table")
                                print(dict_text)
                                print(id_cnt)
                                dict_text = dict_slide[id_cnt]
                                cell.text = dict_text['text']
                                for paragraph in cell.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.size = Pt(10)
                                id_cnt += 1
                            
                            
                # # Kiểm tra nếu shape là placeholder và có chứa text
                # elif shape.is_placeholder:
                #     if hasattr(shape, "text_frame") and shape.text_frame is not None and id_cnt < len(dict_slide):
                #         dict_text = dict_slide[id_cnt]
                #         print(dict_text)
                #         shape.text_frame.clear()
                #         p = shape.text_frame.add_paragraph()
                #         p.text =  dict_text['text'] 
            slide_cnt+=1

    # def extract_text(self, )

if __name__ == "__main__":
    pp_extract = Pineline()
    
    # pass